import sqlite3
import zipfile

import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation


def extract_text(filepath: str, file_type: str) -> str:
    if file_type == "pdf":
        doc = fitz.open(filepath)
        return "\n".join([page.get_text() for page in doc])

    elif file_type == "pptx":
        prs = Presentation(filepath)
        return "\n".join(
            [
                shape.text
                for slide in prs.slides
                for shape in slide.shapes
                if shape.has_text_frame
            ]
        )

    elif file_type == "docx":
        doc = Document(filepath)
        return "\n".join([p.text for p in doc.paragraphs])

    elif file_type in ("txt", "md"):
        with open(filepath, "r") as f:
            return f.read()

    elif file_type == "apkg":
        # Anki deck — zip file containing a SQLite database
        with zipfile.ZipFile(filepath) as z:
            z.extract("collection.anki2", "/tmp/")
        conn = sqlite3.connect("/tmp/collection.anki2")
        notes = conn.execute("SELECT flds FROM notes").fetchall()
        return "\n".join([note[0].replace("\x1f", " | ") for note in notes])

    return ""


def chunk_text(
    text: str, source: str, chunk_size: int = 500, overlap: int = 50
) -> list[dict]:
    words = text.split()
    chunks = []
    i = 0
    while i < len(words):
        chunk = " ".join(words[i : i + chunk_size])
        chunks.append({"source": source, "text": chunk})
        i += chunk_size - overlap
    return chunks
